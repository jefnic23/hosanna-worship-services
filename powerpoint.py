import os
import re

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.util import Inches, Pt

from utils import get_height, get_parts, get_width, grouper, lookahead, get_superscripts

MAX_WIDTH = 565
MAX_HEIGHT = 335


regular = ImageFont.truetype('%SystemRoot%\Fonts\segoeui.ttf', 24)
bold = ImageFont.truetype('%SystemRoot%\Fonts\segoeuib.ttf', 24)
draw = ImageDraw.Draw(Image.new('RGB', (MAX_WIDTH, MAX_HEIGHT)))



class PowerPoint():
    '''Creates a PowerPoint presentation'''

    DEFAULT_FONT = 'Segoe UI'
    DEFAULT_FONTSIZE = 18
    
    def __init__(self, day):
        self._day = day
        self.prs = Presentation()

        self.prs.slide_width = Inches(6)
        self.prs.slide_height = Inches(4)

        self._section_layout = self.prs.slide_layouts[2]
        self._blank_layout = self.prs.slide_layouts[6]

        if os.path.exists(f'services/{day}/image.pptx'):
            self._get_image()

        if os.path.exists(f'services/{day}/hymns.txt'):
            self._hymns = self._load_hymns()

        if os.path.exists('liturgy/confession.txt'):
            self._confession = open('liturgy/confession.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/greeting.txt'):
            self._greeting = open('liturgy/greeting.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/kyrie.txt'):
            self._kyrie = open('liturgy/kyrie.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/creed.txt'):
            self._creed = open('liturgy/creed.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/dialogue.txt'):
            self._dialogue = open('liturgy/dialogue.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/holy-holy-holy.txt'):
            self._hosanna = open('liturgy/holy-holy-holy.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/communion-dialogue.txt'):
            self._communion_dialogue = open('liturgy/communion-dialogue.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/lords-prayer.txt'):
            self._lords_prayer = open('liturgy/lords-prayer.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/communion-hymn.txt'):
            self._communion_hymn = open('liturgy/communion-hymn.txt', 'r', encoding='utf-8').read()

        if os.path.exists(f'services/{day}/prayer.txt'):
            self._prayer = open(f'services/{day}/prayer.txt', 'r', encoding='utf-8').read()

        if os.path.exists(f'services/{day}/first-reading.txt'):
            self._first_reading = open(f'services/{day}/first-reading.txt', 'r', encoding='utf-8').read()

        if os.path.exists(f'services/{day}/psalm.txt'):
            self._psalm = open(f'services/{day}/psalm.txt', 'r', encoding='utf-8').read()

        if os.path.exists(f'services/{day}/second-reading.txt'):
            self._second_reading = open(f'services/{day}/second-reading.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/gospel-acclamation.txt'):
            self._gospel_acclamation = open('liturgy/gospel-acclamation.txt', 'r', encoding='utf-8').read()
        
        if os.path.exists(f'services/{day}/gospel.txt'):
            self._gospel = open(f'services/{day}/gospel.txt', 'r', encoding='utf-8').read()

        if os.path.exists(f'services/{day}/intercession.txt'):
            self._intercession = open(f'services/{day}/intercession.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/prayer-after-communion.txt'):
            self._prayer_after_communion = open('liturgy/prayer-after-communion.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/benediction.txt'):
            self._benediction = open('liturgy/benediction.txt', 'r', encoding='utf-8').read()

        if os.path.exists('liturgy/dismissal.txt'):
            self._dismissal = open('liturgy/dismissal.txt', 'r', encoding='utf-8').read()


    def add_image(self):
        slide = self.prs.slides.add_slide(self._blank_layout)
        left = top = Inches(0)
        slide.shapes.add_picture(f'services/{self._day}/image.jpg', left, top, self.prs.slide_width, self.prs.slide_height)


    def add_confession(self, title_text='Confession and Forgiveness'):
        '''Add the confession to the presentation'''
        pastor = [p.start() for p in re.finditer(r'P:', self._confession)]
        congregation = [c.start() for c in re.finditer(r'C:', self._confession)]

        for part in get_parts(pastor, congregation):
            if part[1] is not None:
                p, c = self._confession[part[0][0]:part[0][1]].strip(), self._confession[part[1][0]:part[1][1]].strip()
            else:
                p, c = self._confession[part[0][0]:part[0][1]].strip(), self._confession[part[0][1]:].strip()
                
            width_formatted_text = []
            for line in p.splitlines():
                width_formatted_text.append(get_width(line, draw, regular))
                
            for line in c.splitlines():
                width_formatted_text.append(get_width(line, draw, bold))

            width_formatted_text = '\n'.join(width_formatted_text)

            slides = get_height(width_formatted_text, draw, regular)

            for slide in slides:
                s = self._add_slide_with_header(title_text)
                content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
                tf = content.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                for line, has_more in lookahead(slide.splitlines()):
                    paragraph = tf.paragraphs[0]
                    self._add_run(paragraph, line, bold=True if line in c else False, has_more=has_more)


    def add_prayer_of_the_day(self, title_text='Prayer of the Day'):
        '''Add the prayer of the day to the presentation'''
        slide = self._add_slide_with_header(title_text)
        content = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        for line, has_more in lookahead(self._prayer.splitlines()):
            paragraph = tf.paragraphs[0]
            paragraph.alignment = PP_ALIGN.JUSTIFY
            self._add_run(paragraph, line, bold=True, has_more=has_more)


    def add_greeting(self, title_text='Greeting'):
        '''Add the greeting to the presentation'''
        pastor = [p.start() for p in re.finditer(r'P:', self._greeting)]
        congregation = [c.start() for c in re.finditer(r'C:', self._greeting)]

        for part in get_parts(pastor, congregation):
            if part[1] is not None:
                p, c = self._greeting[part[0][0]:part[0][1]].strip(), self._greeting[part[1][0]:part[1][1]].strip()
            else:
                p, c = self._greeting[part[0][0]:part[0][1]].strip(), self._greeting[part[0][1]:].strip()
                
            width_formatted_text = []
            for line in p.splitlines():
                width_formatted_text.append(get_width(line, draw, regular))
                
            for line in c.splitlines():
                width_formatted_text.append(get_width(line, draw, bold))

            width_formatted_text = '\n'.join(width_formatted_text)

            slides = get_height(width_formatted_text, draw, regular)

            for slide in slides:
                s = self._add_slide_with_header(title_text)
                content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
                tf = content.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                for line, has_more in lookahead(slide.splitlines()):
                    paragraph = tf.paragraphs[0]
                    self._add_run(paragraph, line, bold=True if line in c else False, has_more=has_more)


    def add_kyrie(self, title_text='Kyrie'):
        '''Add the Kyrie to the presentation'''
        slides = grouper(self._kyrie.split('\n\n'), 2)

        for slide in slides:
            s = self._add_slide_with_header(title_text)
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
            tf = content.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            for lines, has_more in lookahead(slide) if slide else []:
                paragraph = tf.paragraphs[0]
                for line, has_break in lookahead(lines.splitlines()) if lines else []:
                    self._add_run(paragraph, line, bold=True, has_more=has_break)
                if has_more:
                    paragraph.add_line_break()
                    paragraph.add_line_break()


    def add_hymn(self):
        '''Add a hymn to the presentation'''
        slide = self.prs.slides.add_slide(self._blank_layout)
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for line, has_more in lookahead(self._hymns.pop(0)):
            paragraph = tf.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            self._add_run(paragraph, line, bold=True, italic=True if has_more else False, size=24, has_more=has_more)


    def add_reading(self, title_text, text):
        '''Add a reading to the presentation'''
        title = text.splitlines()[0]
        reading = text.splitlines()[1:]

        self.add_title_slide(title)

        width_formatted_text = []
        for line in reading:
            width_formatted_text.append(get_width(line, draw, regular))

        width_formatted_text = '\n'.join(width_formatted_text)

        slides = get_height(width_formatted_text, draw, regular)

        for slide, is_not_last in lookahead(slides):
            s = self._add_slide_with_header(title_text)
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
            tf = content.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            paragraph = tf.paragraphs[0]
            for line, has_more in lookahead(slide.splitlines()):
                superscripts = get_superscripts(line)
                if superscripts:
                    for start, end in superscripts:
                        self._add_run(paragraph, line[:start])
                        self._add_run(paragraph, line[start:end], superscript=True)
                        self._add_run(paragraph, line[end:], has_more=has_more)
                else:
                    if not is_not_last and not has_more:
                        self._add_run(paragraph, line, bold=True)
                    else:
                        self._add_run(paragraph, line, has_more=has_more)


    def add_psalm(self):
        '''Add a psalm to the presentation'''
        text = self._psalm.replace('|', '').replace('- ', '')
        title = text.splitlines()[0]
        psalm = [line.strip() for line in text.splitlines()[1:] if line]

        self.add_title_slide(title)

        p = []
        c = []
        first_lines = []

        width_formatted_text = []
        for i, line in enumerate(psalm):
            if i % 2 == 0:
                new_line = get_width(line, draw, regular)
                first_lines.append(new_line.splitlines()[0])
                for l in new_line.splitlines():
                    p.append(l)
                width_formatted_text.append(new_line)
            else:
                new_line = get_width(line, draw, bold)
                first_lines.append(new_line.splitlines()[0])
                for l in new_line.splitlines():
                    c.append(l)
                width_formatted_text.append(new_line)

        width_formatted_text = '\n'.join(width_formatted_text)

        slides = get_height(width_formatted_text, draw, regular)

        for slide in slides:
            s = self._add_slide_with_header(title)
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
            tf = content.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            paragraph = tf.paragraphs[0]
            for line, has_more in lookahead(slide.splitlines()):
                if line in first_lines:
                    superscripts = get_superscripts(line)
                    if superscripts:
                        for start, end in superscripts:
                            self._add_run(paragraph, line[:start], bold=True if line in c else False)
                            self._add_run(paragraph, line[start:end], bold=True if line in c else False, superscript=True)
                            self._add_run(paragraph, line[end:], bold=True if line in c else False, has_more=has_more)
                    else:
                        self._add_run(paragraph, line, bold=True, has_more=has_more)
                else:
                    self._add_run(paragraph, line, bold=True if line in c else False, has_more=has_more)


    def add_gospel_acclamation(self):
        '''Add the gospel acclamation to the presentation'''
        slide = self._add_slide_with_header('Gospel Acclamation')
        content = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph = tf.paragraphs[0]
        for line, has_more in lookahead(self._gospel_acclamation.splitlines()):
            self._add_run(paragraph, line, bold=True, has_more=has_more)


    def add_gospel(self):
        '''Add the gospel to the presentation'''
        text = self._gospel
        title = text.splitlines()[0].split()[0]
        reading = text.splitlines()[1:]

        slide = self.prs.slides.add_slide(self._blank_layout)
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        self._add_run(p, f'The holy gospel according to {title}.')
        p.add_line_break()
        self._add_run(p, 'Glory to you, O Lord.', bold=True)

        width_formatted_text = []
        for line in reading:
            width_formatted_text.append(get_width(line, draw, regular))

        width_formatted_text = '\n'.join(width_formatted_text)

        slides = get_height(width_formatted_text, draw, regular)

        for slide, is_not_last in lookahead(slides):
            s = self._add_slide_with_header('Gospel')
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
            tf = content.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            paragraph = tf.paragraphs[0]
            for line, has_more in lookahead(slide.splitlines()):
                superscripts = get_superscripts(line)
                if superscripts:
                    for start, end in superscripts:
                        self._add_run(paragraph, line[:start])
                        self._add_run(paragraph, line[start:end], superscript=True)
                        self._add_run(paragraph, line[end:], has_more=has_more)
                else:
                    if not is_not_last and not has_more:
                        self._add_run(paragraph, line, bold=True)
                    else:
                        self._add_run(paragraph, line, has_more=has_more)


    def add_creed(self, title):
        '''Add the creed to the presentation'''
        width_formatted_text = []
        for line in self._creed.splitlines():
            width_formatted_text.append(get_width(line, draw, bold))

        width_formatted_text = '\n'.join(width_formatted_text)

        slides = get_height(width_formatted_text, draw, regular)

        for slide in slides:
            s = self._add_slide_with_header(title)
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
            tf = content.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            p = tf.paragraphs[0]
            for line, has_more in lookahead(slide.splitlines()):
                self._add_run(p, line, bold=True, has_more=has_more)


    def add_intercessions(self):
        '''Add the intercessions to the presentation'''
        slide = self._add_slide_with_header('Prayers of Intercession')
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        self._add_run(p, 'Each petition ends:', italic=True)
        p.add_line_break()
        self._add_run(p, self._intercession.splitlines()[0])
        p.add_line_break()
        self._add_run(p, self._intercession.splitlines()[1], bold=True)


    def add_dialogue(self):
        '''Add the dialogue to the presentation'''
        text = grouper(self._dialogue.splitlines(), 2)

        slide = self._add_slide_with_header('Dialogue')
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        for line, has_more in lookahead(text):
            self._add_run(p, line[0])
            p.add_line_break()
            self._add_run(p, line[1], bold=True, has_more=has_more)
            if has_more:
                p.add_line_break()


    def add_hosanna(self):
        '''Add holy holy holy to the presentation'''
        slide = self._add_slide_with_header('Holy, holy, holy')
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        for line in self._hosanna.splitlines():
            self._add_run(p, line, bold=True)
            p.add_line_break()


    def add_communion_dialogue(self):
        '''Add the communion dialogue to the presentation'''
        slide = self._add_slide_with_header('Communion Dialogue')
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        for line, has_more in lookahead(self._communion_dialogue.splitlines()):
            self._add_run(p, line, bold=False if has_more else True)
            p.add_line_break()


    def add_lords_prayer(self, title):
        '''Add the creed to the presentation'''
        width_formatted_text = []
        for line in self._lords_prayer.splitlines():
            width_formatted_text.append(get_width(line, draw, bold))

        width_formatted_text = '\n'.join(width_formatted_text)

        slides = get_height(width_formatted_text, draw, regular)

        for slide in slides:
            s = self._add_slide_with_header(title)
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
            tf = content.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            p = tf.paragraphs[0]
            for line, has_more in lookahead(slide.splitlines()):
                self._add_run(p, line, bold=True, has_more=has_more)


    def add_communion_hymn(self):
        '''Add the communion hymn to the presentation'''
        slides = grouper(self._communion_hymn.split('\n\n'), 2)

        for slide in slides:
            s = self._add_slide_with_header('Communion Hymn')
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
            tf = content.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            for lines, has_more in lookahead(slide) if slide else []:
                paragraph = tf.paragraphs[0]
                for line, has_break in lookahead(lines.splitlines()) if lines else []:
                    self._add_run(paragraph, line, bold=True, has_more=has_break)
                if has_more:
                    paragraph.add_line_break()
                    paragraph.add_line_break()


    def add_call_and_response(self, title, text):
        '''Add the call and response to the presentation'''
        width_formatted_text = []
        for line in text.splitlines():
            width_formatted_text.append(get_width(line, draw, bold))

        width_formatted_text = '\n'.join(width_formatted_text)

        slides = get_height(width_formatted_text, draw, regular)

        for slide in slides:
            s = self._add_slide_with_header(title)
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(0))
            tf = content.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            p = tf.paragraphs[0]
            for line, has_more in lookahead(slide.splitlines()):
                self._add_run(p, line, bold=True if not has_more else False, has_more=has_more)


    def save(self):
        '''Save the presentation'''
        if not os.path.exists(f'services/{self._day}'):
            os.makedirs(f'services/{self._day}')
        self.prs.save(f'services/{self._day}/service.pptx')


    def _load_hymns(self):
        '''Load the hymns from the hymns.txt file'''
        hymns = open(f'services/{self._day}/hymns.txt', 'r', encoding='utf-8').read()
        return grouper(hymns.splitlines(), 2)


    def _get_layouts(self):
        '''Get the layouts of the presentation'''
        for l in self.prs.slide_layouts:
            print(l.name)


    def _get_image(self):
        '''Add an image to the presentation'''
        prs = Presentation(f'services/{self._day}/image.pptx')
        slide = prs.slides[0]
        shape = slide.shapes[0]
        image = shape.image
        blob, ext = image.blob, image.ext
        with open(f'services/{self._day}/image.{ext}', 'wb') as f:
            f.write(blob)
    
        os.remove(f'services/{self._day}/image.pptx')


    def _add_slide_with_header(self, title_text):
        '''Add a slide with a header'''
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._add_header(slide, title_text)
        return slide
    

    def add_title_slide(self, text):
        '''Add a title slide to the presentation'''
        slide = self.prs.slides.add_slide(self._blank_layout)
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._add_run(p, text, bold=True, size=24)
    

    @staticmethod
    def _get_pastor(text):
        '''Returns the indices of the pastor's lines'''
        return [p.start() for p in re.finditer(r'P:', text)]


    @staticmethod
    def _get_congregation(text):
        '''Returns the indices of the congregation's lines'''
        return [c.start() for c in re.finditer(r'C:', text)]

    
    @staticmethod
    def _add_header(slide, header_text):
        '''Add a header to a slide'''
        header = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(0))
        tf = header.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph = tf.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        PowerPoint._add_run(paragraph, header_text, size=12, bold=True, color=(66, 133, 244))
    
    
    @staticmethod
    def _add_run(paragraph, text, font=DEFAULT_FONT, size=DEFAULT_FONTSIZE, bold=False, italic=False, color=(0, 0, 0), has_more=False, superscript=False):
        '''Add a run to a paragraph'''
        run = paragraph.add_run()
        run.text = text.strip()
        run.font.name = font
        run.font.bold = bold
        run.font.italic = italic
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(*color)
        run.font._element.set('baseline', '30000' if superscript else '0')
        if has_more:
            paragraph.add_line_break()
        
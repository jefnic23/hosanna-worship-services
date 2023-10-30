import os
import re
from datetime import date
from itertools import chain, pairwise
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from PIL.ImageFont import FreeTypeFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.slide import Slide
from pptx.text.text import _Paragraph
from pptx.util import Inches, Pt

from models.hymn import Hymn
from models.reading import Reading
from services.settings import Settings
from services.utils import get_superscripts, lookahead, split_formatted_text


class PowerPoint:
    """
    A class for creating a PowerPoint presentation.
    """
    DEFAULT_FONT: str = 'Segoe UI'
    DEFAULT_FONTSIZE: int = 18
    MAX_WIDTH: int = 565
    MAX_HEIGHT: int = 335
    
    REGULAR: FreeTypeFont = ImageFont.truetype('%SystemRoot%\Fonts\segoeui.ttf', 24)
    BOLD: FreeTypeFont = ImageFont.truetype('%SystemRoot%\Fonts\segoeuib.ttf', 24)
    ITALIC: FreeTypeFont = ImageFont.truetype('%SystemRoot%\Fonts\segoeuii.ttf', 24)
    DRAW: ImageDraw.ImageDraw = ImageDraw.Draw(Image.new('RGB', (MAX_WIDTH, MAX_HEIGHT)))
    
    def __init__(self, settings: Settings, title: str):
        self.day: date = date.today()
        self.title: str = title
        self.prs = Presentation()
        self.prs.slide_width = Inches(6)
        self.prs.slide_height = Inches(4)
        self._section_layout = self.prs.slide_layouts[2]
        self._blank_layout = self.prs.slide_layouts[6]
        self._path = f'{settings.LOCAL_DIR}/services'

        # if os.path.exists(f'{path}/{day}/hymns.txt'):
        #     self._hymns = self._load_hymns()


    def add_title_slide(self, text: str) -> None:
        '''Add a title slide to the presentation/'''
        slide = self.prs.slides.add_slide(self._blank_layout)
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._add_run(p, text, bold=True, size=24)


    def add_image(self) -> None:
        '''Add an image to the presentation.'''
        slide = self.prs.slides.add_slide(self._blank_layout)
        left = top = Inches(0)
        slide.shapes.add_picture(
            f'{self._path}/{self.day}/image.jpg', 
            left, 
            top,
            self.prs.slide_width, 
            self.prs.slide_height
        )
        # TODO: delete image after adding to presentation


    def add_hymn(self, hymn: Hymn) -> None:
        '''Add a hymn to the presentation.'''
        slide = self.prs.slides.add_slide(self._blank_layout)
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for i, value in enumerate([hymn.title, hymn.number]):
            paragraph = tf.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            self._add_run(
                paragraph, 
                value, 
                bold=True, 
                italic=True if i == 0 else False, 
                size=24, 
                has_more=True if i == 0 else False
            )


    def add_gospel_title(self, title: str) -> None:
        '''Add the title of the Gospel reading to the presentation.'''
        slide = self.prs.slides.add_slide(self._blank_layout)
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        self._add_run(p, f'The holy gospel according to {title.split()[0]}.')
        p.add_line_break()
        self._add_run(p, 'Glory to you, O Lord.', bold=True)


    def add_intercessions(
        self, 
        call: str,
        response: str
    ) -> None:
        '''Add the intercessions to the presentation.'''
        slide = self._add_slide_with_header('Prayers of Intercession')
        content = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(4))
        tf = content.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        self._add_run(p, 'Each petition ends:', italic=True)
        p.add_line_break()
        self._add_run(p, call)
        p.add_line_break()
        self._add_run(p, response, bold=True)


    def add_rich_text(
            self, 
            title: str, 
            text: str, 
            anchor: str = '',
            spoken: bool = False,
            draw: ImageDraw.ImageDraw = DRAW,
            regular: FreeTypeFont = REGULAR,
            bold: FreeTypeFont = BOLD,
            italic: FreeTypeFont = ITALIC
        ) -> None:
        """Adds rich text to the presentation.

        Args:
            title (str): the title of the slide. Gets added to the top right.
            text (str): the body of the slide.
            anchor (str, optional): _description_. Defaults to ''.
            spoken (bool, optional): wraps all text in bold tags if True. Defaults to False.
            draw (ImageDraw.ImageDraw, optional): _description_. Defaults to DRAW.
            regular (FreeTypeFont, optional): _description_. Defaults to REGULAR.
            bold (FreeTypeFont, optional): _description_. Defaults to BOLD.
            italic (FreeTypeFont, optional): _description_. Defaults to ITALIC.
        """
        if spoken:
            text = '\n'.join([f'<b>{line}</b>' for line in text.splitlines()])
        superscripts = re.findall(r'<sup>(.*?)</sup>', text)
        regular_text, bold_text, italic_text = split_formatted_text(text.replace('<sup>', '').replace('</sup>', ''))
        bold_lines = [line[1] for line in bold_text]
        italic_lines = [line[1] for line in italic_text]
        lines = [
            line[1] for line in sorted(regular_text + bold_text + italic_text, key=lambda x: x[0])
        ]
        width_formatted_text = []
        bold_formatted_text = []
        italic_formatted_text = []
        for line in lines:
            if line in bold_lines:
                bold_line = PowerPoint.get_width(line, draw, bold)
                width_formatted_text.append(bold_line)
                for l in bold_line.splitlines():
                    bold_formatted_text.append(l)
            elif line in italic_lines:
                italic_line = PowerPoint.get_width(line, draw, italic)
                width_formatted_text.append(italic_line)
                for l in italic_line.splitlines():
                    italic_formatted_text.append(l)
            else:
                width_formatted_text.append(PowerPoint.get_width(line, draw, regular))

        width_formatted_text = '\n'.join(width_formatted_text)

        slides = PowerPoint.get_height(width_formatted_text, draw, regular)

        for slide, is_not_last in lookahead(slides):
            s = self._add_slide_with_header(title)
            content = s.shapes.add_textbox(Inches(0), Inches(0.5), Inches(6), Inches(3))
            tf = content.text_frame
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = None if not anchor else PowerPoint._anchor_map(anchor)
            paragraph = tf.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            for line, has_more in lookahead(slide.splitlines()):
                sups = get_superscripts(superscripts, line)
                if len(sups) > 0:
                    index = pairwise(list(chain(*[
                        [0], 
                        *[[s, e] for s, e in sups], 
                        [len(line)]
                    ])))
                    for start, end in index:
                        self._add_run(
                            paragraph, 
                            line[start:end], 
                            superscript=True if (start, end) in sups else False,
                            bold=True if line in bold_formatted_text else False,
                            italic=True if line in italic_formatted_text else False,
                            color=(86, 86, 86) if (start, end) in sups else (0, 0, 0)
                        )
                    if has_more:
                        paragraph.add_line_break()
                    for _ in range(len(sups)):
                        superscripts.pop(0)
                else:
                    if not is_not_last and not has_more:
                        self._add_run(
                            paragraph, 
                            line, 
                            bold=True if line in bold_formatted_text else False,
                            italic=True if line in italic_formatted_text else False
                        )
                    else:
                        self._add_run(
                            paragraph, 
                            line, 
                            has_more=has_more, 
                            bold=True if line in bold_formatted_text else False,
                            italic=True if line in italic_formatted_text else False
                        )


    # TODO: add image file name to method signature
    def convert_image(self) -> None:
        """Images downloaded from Sundays and Seasons are in a .pptx format. 
        This method converts the image to a .jpg that can be added to the powerpoint.
        """
        prs = Presentation(f'{self._path}/{self.day}/image.pptx')
        slide = prs.slides[0]
        shape = slide.shapes[0]
        image = shape.image
        blob, ext = image.blob, image.ext
        with open(f'{self._path}/{self.day}/image.{ext}', 'wb') as f:
            f.write(blob)
    
        os.remove(f'{self._path}/{self.day}/image.pptx')
    

    def save(self) -> None:
        '''Save the presentation.'''
        if not os.path.exists(f'{self._path}/{self.day}'):
            os.makedirs(f'{self._path}/{self.day}')
        self.prs.save(f'{self._path}/{self.day}/{self.title}.pptx')


    #region Private Methods

    def _get_layouts(self) -> None:
        '''Get the layouts of the presentation.'''
        for layout in self.prs.slide_layouts:
            print(layout.name)


    def _add_slide_with_header(self, title_text: str) -> Slide:
        '''Add a slide with a header'''
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._add_header(slide, title_text)
        return slide
    
    #endregion
    
    #region Static Methods

    @staticmethod
    def _get_pastor(text: str) -> list[int]:
        '''Returns the indices of the pastor's lines'''
        return [p.start() for p in re.finditer(r'P:', text)]


    @staticmethod
    def _get_congregation(text: str) -> list[int]:
        '''Returns the indices of the congregation's lines'''
        return [c.start() for c in re.finditer(r'C:', text)]

    
    @staticmethod
    def _add_header(slide: Slide, header_text: str) -> None:
        '''Add a header to a slide'''
        header = slide.shapes.add_textbox(Inches(0), Inches(0.1), Inches(6), Inches(0))
        tf = header.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        paragraph = tf.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        PowerPoint._add_run(
            paragraph, 
            header_text, 
            size=12, 
            bold=True, 
            color=(66, 133, 244)
        )
    
    
    @staticmethod
    def _add_run(
        paragraph: _Paragraph, 
        text: str, 
        font: str = DEFAULT_FONT, 
        size: int = DEFAULT_FONTSIZE, 
        bold: bool = False, 
        italic: bool = False, 
        color: tuple = (0, 0, 0), 
        has_more: bool = False, 
        superscript: bool = False
    ) -> None:
        '''Add a run to a paragraph'''
        run = paragraph.add_run()
        run.text = text
        run.font.name = font
        run.font.bold = bold
        run.font.italic = italic
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(*color)
        run.font._element.set('baseline', '30000' if superscript else '0')
        if has_more:
            paragraph.add_line_break()


    @staticmethod
    def check_size(
        line: str,
        draw: ImageDraw.ImageDraw,
        font: FreeTypeFont
    ) -> dict[str, int]:
        '''Checks the size of a line of text.'''
        size = draw.multiline_textbbox((0, 0), line, font)
        return {'width': size[2], 'height': size[3]}


    @staticmethod
    def get_width(
        line: str, 
        draw: ImageDraw.ImageDraw, 
        font: FreeTypeFont, 
        max_width: int = MAX_WIDTH
    ):
        '''Gets the width of a line of text and splits it if it's too long.'''
        if PowerPoint.check_size(line, draw, font)['width'] < max_width:
            return line
        else:
            lines = []
            for word in line.split():
                width = PowerPoint.check_size(
                    ' '.join(lines[-1:] + [word]), 
                    draw, 
                    font
                )['width']
                if width < max_width:
                    lines[-1:] = [' '.join(lines[-1:] + [word])]
                else:
                    lines += [word]

            return '\n'.join(lines)
        

    @staticmethod
    def get_height(
        lines: str, 
        draw: ImageDraw.ImageDraw, 
        font: FreeTypeFont, 
        max_height: int = MAX_HEIGHT
    ) -> list[str]:
        '''Gets the height of a line of text and splits it if it's too long.'''
        if PowerPoint.check_size(lines, draw, font)['height'] < max_height:
            return [lines]

        regex = re.search(r'<div>(.*?)</div>', lines, re.MULTILINE | re.DOTALL)
        group = regex.group(1).strip().splitlines() if regex else None
        slides = []
        for line in lines.splitlines():
            if re.match(r'<br>', line):
                # If the line is a break, insert a new slide
                slides += ['']
                continue

            if re.match(r'<div>', line) or re.match(r'</div>', line):
                continue

            height = PowerPoint.check_size(
                '\n'.join(slides[-1:] + [line]), 
                draw, 
                font
            )['height']
            
            if regex is not None and line == group[0]:
                # check if the call and response will fit on the current slide
                height = PowerPoint.check_size(
                    '\n'.join(slides[-1:] + [line] + [group[1]]),
                    draw,
                    font
                )['height']
                if height > max_height:
                    slides += [line]
                else:
                    slides[-1:] = ['\n'.join(slides[-1:] + [line])]
                continue

            if height < max_height:
                slides[-1:] = ['\n'.join(slides[-1:] + [line]).lstrip()]
            else:
                slides += [line]
        return slides
        

    @staticmethod
    def get_slides(
        lines: str, 
        draw: ImageDraw.ImageDraw, 
        regular: FreeTypeFont,
    ) -> list[str]:
        '''Transforms a string into a list of strings that fit on a slide.'''
        width_formatted_text = []
        for line in lines.splitlines():
            width_formatted_text.append(PowerPoint.get_width(line, draw, regular))
        width_formatted_text = '\n'.join(width_formatted_text)

        return PowerPoint.get_height(width_formatted_text, draw, regular)
    

    @staticmethod
    def _anchor_map(anchor: str) -> MSO_VERTICAL_ANCHOR:
        '''Returns the vertical anchor position of a text box.'''
        return {
            'top': MSO_VERTICAL_ANCHOR.TOP, # type: ignore
            'middle': MSO_VERTICAL_ANCHOR.MIDDLE, # type: ignore
            'bottom': MSO_VERTICAL_ANCHOR.BOTTOM # type: ignore
        }[anchor]
            
    #endregion
